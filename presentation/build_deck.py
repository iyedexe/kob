"""Build the kob presentation: presentation/kob.pptx (16:9).

Reproducible, data-driven deck — motivation, a 30-second theory intro, architecture,
key findings, two performance charts (from docs/PERFORMANCE_REPORT.md), and the demo.

Build:
    uv run --with python-pptx --with matplotlib python presentation/build_deck.py
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
ASSETS.mkdir(exist_ok=True)

# ---- palette ----------------------------------------------------------------
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
NAVY  = RGBColor(0x26, 0x46, 0x53)
INK   = RGBColor(0x1D, 0x2B, 0x3A)
SAND  = RGBColor(0xE9, 0xC4, 0x66)
SAND2 = RGBColor(0xF4, 0xA2, 0x61)
ORANGE= RGBColor(0xE7, 0x6F, 0x51)
GRAY  = RGBColor(0x8A, 0x94, 0x9E)
LIGHT = RGBColor(0xF2, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG= RGBColor(0x1E, 0x2A, 0x36)
CODEFG= RGBColor(0xE6, 0xED, 0xF3)

HEXES = {"TEAL": "#2a9d8f", "NAVY": "#264653", "SAND": "#e9c466",
         "SAND2": "#f4a261", "ORANGE": "#e76f51"}

# ============================================================================
# Charts
# ============================================================================
def _ms(v):
    return f"{v:.0f} ms" if v < 1000 else f"{v/1000:.1f} s"

def chart_latency(path):
    labels = ["Arrow Flight", "Arrow IPC / HTTP", "Protobuf (columnar)", "Protobuf (row)", "REST / JSON"]
    vals   = [66.9, 71.5, 3741.2, 6475.6, 117128.9]   # ms, 1.75M-row pull (medium dataset)
    colors = [HEXES["TEAL"], HEXES["NAVY"], HEXES["SAND"], HEXES["SAND2"], HEXES["ORANGE"]]
    fig, ax = plt.subplots(figsize=(8.8, 4.0), dpi=200)
    bars = ax.barh(labels[::-1], vals[::-1], color=colors[::-1], height=0.62)
    ax.set_xscale("log")
    ax.set_xlim(40, 500000)
    ax.set_xlabel("end-to-end client latency (log scale) — lower is better", fontsize=9)
    for b, v in zip(bars, vals[::-1]):
        ax.text(v * 1.18, b.get_y() + b.get_height() / 2, _ms(v), va="center", fontsize=9.5, weight="bold")
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=9.5)
    ax.grid(axis="x", color="#ececec"); ax.set_axisbelow(True)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)

def chart_speedup(path):
    labels = ["Arrow\nFlight", "Arrow IPC\nover HTTP", "Protobuf\ncolumnar", "Protobuf\nrow"]
    vals   = [1751, 1019, 31, 18]
    colors = [HEXES["TEAL"], HEXES["NAVY"], HEXES["SAND"], HEXES["SAND2"]]
    fig, ax = plt.subplots(figsize=(8.6, 3.9), dpi=200)
    bars = ax.bar(labels, vals, color=colors, width=0.62)
    ax.set_yscale("log"); ax.set_ylim(1, 4000)
    ax.set_ylabel("× faster than REST/JSON (log)", fontsize=9)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v * 1.12, f"{v}×", ha="center", fontsize=11, weight="bold")
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(labelsize=9.5)
    ax.grid(axis="y", color="#ececec"); ax.set_axisbelow(True)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)

LAT = ASSETS / "latency.png"
SPD = ASSETS / "speedup.png"
chart_latency(LAT)
chart_speedup(SPD)

# ============================================================================
# Deck helpers
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
MX = Inches(0.75)
CW = Inches(11.8)

def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background(); shape.shadow.inherit = False
    return shape

def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tf

def render(tf, items):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = it.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(it.get("after", 6))
        p.space_before = Pt(it.get("before", 0))
        p.line_spacing = it.get("ls", 1.05)
        indent = "      " * it.get("level", 0)
        glyph = {"dot": "●  ", "dash": "–  ", None: ""}[it.get("bullet")]
        for j, seg in enumerate(it["runs"]) if "runs" in it else enumerate([it]):
            r = p.add_run()
            r.text = (indent + glyph if j == 0 else "") + seg["text"]
            r.font.size = Pt(seg.get("size", it.get("size", 18)))
            r.font.bold = seg.get("bold", it.get("bold", False))
            r.font.italic = seg.get("italic", it.get("italic", False))
            r.font.color.rgb = seg.get("color", it.get("color", INK))
            r.font.name = seg.get("font", it.get("font", "Calibri"))

def slide(idx, title=None, kicker=None):
    s = prs.slides.add_slide(BLANK)
    fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), Inches(7.5)), TEAL)  # brand bar
    if title:
        if kicker:
            render(tbox(s, MX, Inches(0.36), CW, Inches(0.3)),
                   [{"text": kicker.upper(), "size": 12, "bold": True, "color": TEAL}])
        render(tbox(s, MX, Inches(0.62), CW, Inches(0.7)),
               [{"text": title, "size": 30, "bold": True, "color": NAVY}])
        fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MX, Inches(1.34), Inches(2.1), Pt(3)), TEAL)
    render(tbox(s, MX, Inches(7.02), CW, Inches(0.32)),
           [{"runs": [{"text": "kob", "bold": True, "color": TEAL, "size": 10},
                      {"text": "  ·  self-discovering Parquet → Apache Arrow", "color": GRAY, "size": 10},
                      {"text": f"          {idx} / 9", "color": GRAY, "size": 10}]}])
    return s

def code_box(s, x, y, w, h, lines, size=12.5):
    box = fill(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h), CODEBG)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.14)
    render(tf, [{"text": ln, "size": size, "color": CODEFG, "font": "Consolas", "after": 3, "ls": 1.1}
                for ln in lines])
    return box

def card(s, x, y, w, h, head, body, accent=TEAL):
    fill(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h), LIGHT)
    fill(s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.09), h), accent)
    tf = tbox(s, x + Inches(0.28), y + Inches(0.16), w - Inches(0.4), h - Inches(0.3))
    render(tf, [{"text": head, "size": 16, "bold": True, "color": NAVY, "after": 4},
                {"text": body, "size": 12.5, "color": INK, "ls": 1.08}])

def table(s, x, y, w, rows, col_w, head_fill=NAVY, row_h=Inches(0.42)):
    n, m = len(rows), len(rows[0])
    gt = s.shapes.add_table(n, m, x, y, w, row_h * n).table
    for j, cw in enumerate(col_w):
        gt.columns[j].width = cw
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.12); c.margin_top = Inches(0.04); c.margin_bottom = Inches(0.04)
            c.fill.solid(); c.fill.fore_color.rgb = head_fill if i == 0 else (WHITE if i % 2 else LIGHT)
            p = c.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(13 if i == 0 else 12.5)
            r.font.bold = i == 0
            r.font.color.rgb = WHITE if i == 0 else INK
            r.font.name = "Calibri"
    return gt

# ============================================================================
# Slide 1 — Title
# ============================================================================
s = prs.slides.add_slide(BLANK)
fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)), NAVY)
fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.55), Inches(13.333), Inches(0.10)), TEAL)
render(tbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.6)),
       [{"text": "kob", "size": 96, "bold": True, "color": WHITE}])
render(tbox(s, Inches(0.95), Inches(3.4), Inches(11.5), Inches(1.0)),
       [{"text": "A tiny, self-discovering Parquet query server — served as Apache Arrow, fast.",
         "size": 26, "color": RGBColor(0xCF, 0xE3, 0xDF)}])
render(tbox(s, Inches(0.95), Inches(4.25), Inches(11.5), Inches(0.6)),
       [{"text": "Drop Parquet in a folder → a discoverable, filterable, blazing-fast Arrow API.",
         "size": 16, "italic": True, "color": GRAY}])
render(tbox(s, Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.6)),
       [{"runs": [{"text": "Apache Arrow Flight   ·   DuckDB   ·   FastAPI / Swagger          ",
                   "size": 13, "color": RGBColor(0xCF, 0xE3, 0xDF)},
                  {"text": "github.com/iyedexe/kob", "size": 13, "bold": True, "color": TEAL}]}])

# ============================================================================
# Slide 2 — Motivation
# ============================================================================
s = slide(2, "Serving a folder of Parquet shouldn't be hard — or slow", kicker="Motivation")
render(tbox(s, MX, Inches(1.55), CW, Inches(0.7)),
       [{"text": "A producer microservice writes Parquet to a folder. You must serve it, read-only, "
                 "to Python / C# / C++ / Excel clients. Every usual option is overkill or slow:",
         "size": 16, "color": INK, "ls": 1.1}])
card(s, MX, Inches(2.5), Inches(3.75), Inches(1.9), "A warehouse / broker",
     "Heavy infrastructure in front of what is already a queryable, columnar file format.", NAVY)
card(s, Inches(4.78), Inches(2.5), Inches(3.75), Inches(1.9), "A hand-written API",
     "Declare a schema + catalog in code, then redeploy on every new dataset, partition or column.", SAND2)
card(s, Inches(8.81), Inches(2.5), Inches(3.75), Inches(1.9), "A JSON / Protobuf endpoint",
     "Pays the serialization tax — Arrow's FAQ puts (de)serialization at 80–90% of compute cost.", ORANGE)
render(tbox(s, MX, Inches(4.85), CW, Inches(1.4)),
       [{"runs": [{"text": "What we actually want:  ", "size": 17, "bold": True, "color": NAVY},
                  {"text": "self-discovering", "size": 17, "bold": True, "color": TEAL},
                  {"text": " (no hardcoded schema)  ·  ", "size": 17, "color": INK},
                  {"text": "minimal", "size": 17, "bold": True, "color": TEAL},
                  {"text": " (tiny footprint)  ·  ", "size": 17, "color": INK},
                  {"text": "blazing fast", "size": 17, "bold": True, "color": TEAL},
                  {"text": " (no serialization tax).", "size": 17, "color": INK}], "after": 8}])

# ============================================================================
# Slide 3 — What kob is
# ============================================================================
s = slide(3, "What kob is", kicker="Overview")
render(tbox(s, MX, Inches(1.55), CW, Inches(0.7)),
       [{"text": "Point it at a directory; it discovers your datasets, tells clients what they can filter "
                 "on per folder and per column, and streams results as Apache Arrow. ~750 lines of core server.",
         "size": 16, "color": INK, "ls": 1.1}])
cards = [("Tiny", "Four dependencies — DuckDB, PyArrow, FastAPI, uvicorn. No DB, no broker, no config, no schema.", TEAL),
         ("Self-configuring", "Drop new files / partitions / datasets — they appear live. Discovery reads folder names + one file's metadata; never scans data.", NAVY),
         ("Fast", "DuckDB pushes partition/predicate/projection pruning into the scan; results stream as Arrow over Flight — nothing to deserialize.", SAND2),
         ("Safe", "Read-only. Dataset names sandboxed to the root; columns/ops validated; every value a bound parameter (no SQL injection).", ORANGE)]
xs = [MX, Inches(6.7)]
ys = [Inches(2.55), Inches(4.55)]
for k, (h, b, a) in enumerate(cards):
    card(s, xs[k % 2], ys[k // 2], Inches(5.85), Inches(1.8), h, b, a)

# ============================================================================
# Slide 4 — Theory: why Arrow
# ============================================================================
s = slide(4, "Why Apache Arrow — the 30-second theory", kicker="Theory")
render(tbox(s, MX, Inches(1.55), Inches(6.3), Inches(4.6)),
       [{"text": "Row vs. columnar is the whole story:", "size": 16, "bold": True, "color": NAVY, "after": 10},
        {"bullet": "dot", "text": "JSON & Protobuf are row-oriented. The sender encodes every value "
                                  "field-by-field; the receiver decodes every value and rebuilds columns.", "size": 14.5, "after": 8, "ls": 1.1},
        {"bullet": "dot", "text": "Cost scales with rows × cols — and grows super-linearly with result size.", "size": 14.5, "after": 8, "ls": 1.1},
        {"bullet": "dot", "text": "Arrow's columnar wire layout IS its in-memory layout. Reading a result is a "
                                  "header parse + buffer-pointer setup — essentially nothing to deserialize.", "size": 14.5, "after": 8, "ls": 1.1},
        {"bullet": "dot", "text": "Arrow Flight uses Protobuf only for its tiny control envelope (tickets), "
                                  "never for the bulk data.", "size": 14.5, "ls": 1.1}])
table(s, Inches(7.35), Inches(2.0), Inches(5.2),
      [["Format", "Layout", "Verdict"],
       ["JSON", "row, text", "slowest"],
       ["Protobuf", "row, binary", "binary ≠ fast"],
       ["Arrow IPC", "columnar", "near-zero decode"],
       ["Arrow Flight", "columnar + gRPC", "the default"]],
      col_w=[Inches(1.9), Inches(1.9), Inches(1.4)])
render(tbox(s, Inches(7.35), Inches(4.95), Inches(5.2), Inches(1.2)),
       [{"text": "“(De)serialization can represent 80–90% of computing costs.”", "size": 13.5, "italic": True, "color": NAVY, "after": 3},
        {"text": "— Apache Arrow FAQ. That is exactly the cost Arrow removes.", "size": 12, "color": GRAY}])

# ============================================================================
# Slide 5 — Architecture
# ============================================================================
s = slide(5, "Architecture: Flight + Swagger, one command", kicker="How it works")
render(tbox(s, MX, Inches(1.5), Inches(6.2), Inches(4.5)),
       [{"text": "One straight line, folder → Arrow stream:", "size": 15, "bold": True, "color": NAVY, "after": 8},
        {"bullet": "dot", "text": "core — discover (catalog), validate the query (contract), execute on DuckDB → Arrow (engine).", "size": 13.5, "after": 6, "ls": 1.08},
        {"bullet": "dot", "text": "server — the kob entry point: Flight + a Swagger HTTP control plane, one process, one engine.", "size": 13.5, "after": 6, "ls": 1.08},
        {"bullet": "dot", "text": "demos — Arrow-over-HTTP / JSON / Protobuf servers, only to demonstrate the alternatives.", "size": 13.5, "after": 6, "ls": 1.08},
        {"bullet": "dot", "text": "tools — sample-data generator, reference client, benchmark.", "size": 13.5, "ls": 1.08}])
table(s, Inches(7.2), Inches(1.7), Inches(5.35),
      [["", "Arrow Flight", "Swagger HTTP"],
       ["for", "machines, fast", "humans, discovery"],
       ["port", "8815 (gRPC)", "8000 (/docs)"],
       ["payload", "Arrow, zero-decode", "JSON / CSV"]],
      col_w=[Inches(1.25), Inches(2.05), Inches(2.05)])
code_box(s, Inches(7.2), Inches(4.45), Inches(5.35), Inches(1.5),
         ["$ uv run kob", "  data  (fast)    Flight   grpc://…:8815",
          "  docs  (Swagger) HTTP     http://…:8000/docs"], size=12.5)

# ============================================================================
# Slide 6 — Key findings
# ============================================================================
s = slide(6, "Key findings", kicker="Results")
render(tbox(s, MX, Inches(1.55), Inches(6.0), Inches(4.6)),
       [{"bullet": "dot", "runs": [{"text": "Arrow wins by 1–3 orders of magnitude", "size": 15, "bold": True, "color": NAVY},
                                   {"text": " — 150–1750× faster than JSON, 5–56× faster than the best Protobuf.", "size": 15, "color": INK}], "after": 10, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "Binary ≠ fast.", "size": 15, "bold": True, "color": NAVY},
                                   {"text": " Protobuf is binary and still 5–56× slower — it must parse every value.", "size": 15, "color": INK}], "after": 10, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "The gap grows with result size", "size": 15, "bold": True, "color": NAVY},
                                   {"text": " — ~150× on 29k rows → ~1750× on 1.75M rows.", "size": 15, "color": INK}], "after": 10, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "JSON is a bandwidth hog too", "size": 15, "bold": True, "color": NAVY},
                                   {"text": " — 8× heavier on the wire than Arrow + zstd.", "size": 15, "color": INK}], "after": 10, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "Same engine throughout.", "size": 15, "bold": True, "color": NAVY},
                                   {"text": " Only the wire format changes — so the wire is the cause.", "size": 15, "color": INK}], "ls": 1.1}])
s.shapes.add_picture(str(SPD), Inches(7.0), Inches(1.7), width=Inches(5.6))
render(tbox(s, Inches(7.0), Inches(6.35), Inches(5.6), Inches(0.4)),
       [{"text": "Speedup vs REST/JSON, 1.75M-row pull (medium dataset).", "size": 11, "italic": True, "color": GRAY, "align": PP_ALIGN.CENTER}])

# ============================================================================
# Slide 7 — Performance comparison
# ============================================================================
s = slide(7, "Performance: same query, five wire formats", kicker="Benchmark")
s.shapes.add_picture(str(LAT), Inches(1.35), Inches(1.65), width=Inches(8.4))
render(tbox(s, Inches(10.0), Inches(1.9), Inches(2.6), Inches(4.6)),
       [{"text": "1.75M rows", "size": 15, "bold": True, "color": NAVY, "after": 2},
        {"text": "one query, one engine", "size": 11.5, "color": GRAY, "after": 14},
        {"runs": [{"text": "Flight  ", "size": 14, "bold": True, "color": TEAL}, {"text": "67 ms", "size": 14, "color": INK}], "after": 6},
        {"runs": [{"text": "JSON  ", "size": 14, "bold": True, "color": ORANGE}, {"text": "117 s", "size": 14, "color": INK}], "after": 14},
        {"text": "Same result, ~1750× apart.", "size": 13, "italic": True, "color": NAVY, "ls": 1.15, "after": 10},
        {"text": "Pushdown happens before serialization — and compounds with picking a fast wire.", "size": 11.5, "color": GRAY, "ls": 1.15}])
render(tbox(s, Inches(1.35), Inches(6.55), Inches(8.4), Inches(0.4)),
       [{"text": "Apple M4, localhost loopback (understates Arrow's win on a real network). Source: docs/PERFORMANCE_REPORT.md",
         "size": 10.5, "italic": True, "color": GRAY, "align": PP_ALIGN.CENTER}])

# ============================================================================
# Slide 8 — Demo
# ============================================================================
s = slide(8, "Demo: one query, every transport", kicker="See it")
render(tbox(s, MX, Inches(1.5), Inches(6.1), Inches(4.6)),
       [{"text": "A runnable notebook spins up all four transports, runs the SAME query over each, "
                 "checks they return identical data, and times them.", "size": 15, "color": INK, "after": 12, "ls": 1.12},
        {"bullet": "dot", "text": "Arrow Flight  ·  Arrow-over-HTTP  ·  REST/JSON  ·  gRPC/Protobuf", "size": 14, "bold": True, "color": NAVY, "after": 8, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "Result:  ", "size": 14, "color": INK},
                                   {"text": "all four return the identical 14,006 rows", "size": 14, "bold": True, "color": TEAL},
                                   {"text": " — Flight fastest by far.", "size": 14, "color": INK}], "after": 8, "ls": 1.1},
        {"bullet": "dot", "text": "Discovery, filtering and serialization are all kob's job — zero schema config.", "size": 14, "color": INK, "ls": 1.1}])
code_box(s, Inches(7.0), Inches(1.7), Inches(5.55), Inches(1.35),
         ["# install + launch the demo", "uv sync --extra demo", "make notebook"], size=13)
code_box(s, Inches(7.0), Inches(3.35), Inches(5.55), Inches(2.7),
         ['{', '  "dataset": "optionmetrics",',
          '  "columns": ["date","strike","impl_vol","delta"],',
          '  "filters": [',
          '    {"column":"underlying","op":"=","value":"AAPL"},',
          '    {"column":"delta","op":">=","value":0.4}',
          '  ]',
          '}'], size=11.5)
render(tbox(s, Inches(7.0), Inches(6.15), Inches(5.55), Inches(0.4)),
       [{"text": "notebooks/kob_transports_demo.ipynb", "size": 11, "italic": True, "color": GRAY, "align": PP_ALIGN.CENTER}])

# ============================================================================
# Slide 9 — Takeaways
# ============================================================================
s = prs.slides.add_slide(BLANK)
fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)), NAVY)
fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.7), Inches(13.333), Inches(0.06)), TEAL)
render(tbox(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9)),
       [{"text": "Takeaways", "size": 34, "bold": True, "color": WHITE}])
render(tbox(s, Inches(0.95), Inches(2.0), Inches(11.6), Inches(3.6)),
       [{"bullet": "dot", "runs": [{"text": "Use Arrow on the wire, default to Flight. ", "size": 18, "bold": True, "color": TEAL},
                                   {"text": "Not JSON (150–1750× slower, 8× heavier), not Protobuf (5–56× slower).", "size": 18, "color": RGBColor(0xDF,0xEA,0xE8)}], "after": 14, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "Self-discovering beats configured. ", "size": 18, "bold": True, "color": TEAL},
                                   {"text": "Point kob at a folder; new data appears live, no restart, no schema.", "size": 18, "color": RGBColor(0xDF,0xEA,0xE8)}], "after": 14, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "Minimal & safe by construction. ", "size": 18, "bold": True, "color": TEAL},
                                   {"text": "Four deps, ~750-line core, read-only, parameterised, sandboxed.", "size": 18, "color": RGBColor(0xDF,0xEA,0xE8)}], "after": 14, "ls": 1.1},
        {"bullet": "dot", "runs": [{"text": "It's the wire, not the engine. ", "size": 18, "bold": True, "color": TEAL},
                                   {"text": "Hold DuckDB constant and the format alone moves latency by 1750×.", "size": 18, "color": RGBColor(0xDF,0xEA,0xE8)}], "ls": 1.1}])
render(tbox(s, Inches(0.95), Inches(6.1), Inches(11.6), Inches(0.8)),
       [{"runs": [{"text": "README.md  ·  docs/DESIGN.md  ·  docs/PERFORMANCE_REPORT.md          ", "size": 13, "color": GRAY},
                  {"text": "github.com/iyedexe/kob", "size": 13, "bold": True, "color": TEAL}]}])

OUT = HERE / "kob.pptx"
prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
print(f"charts: {LAT.name}, {SPD.name}")
